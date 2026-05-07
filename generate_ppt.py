from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Colors
DARK_BG   = RGBColor(0x0A, 0x0F, 0x1E)
BLUE      = RGBColor(0x3B, 0x82, 0xF6)
VIOLET    = RGBColor(0x8B, 0x5C, 0xF6)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x94, 0xA3, 0xB8)
EMERALD   = RGBColor(0x10, 0xB9, 0x81)
AMBER     = RGBColor(0xF5, 0x9E, 0x0B)
RED       = RGBColor(0xEF, 0x44, 0x44)
CYAN      = RGBColor(0x06, 0xB6, 0xD4)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(blank_layout)

def bg(slide, color=DARK_BG):
    shape = slide.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def txbox(slide, text, left, top, width, height,
          size=24, bold=False, color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tb

def rect(slide, left, top, width, height, fill_color, alpha_sim=None, line_color=None):
    shape = slide.shapes.add_shape(1,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

def accent_bar(slide, color=BLUE):
    rect(slide, 0, 0, 0.08, 7.5, color)

def bullet_slide(slide, title, items, icon_color=BLUE):
    accent_bar(slide, icon_color)
    txbox(slide, title, 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
    rect(slide, 0.4, 1.15, 2.5, 0.04, icon_color)
    y = 1.5
    for item in items:
        rect(slide, 0.55, y + 0.12, 0.12, 0.12, icon_color)
        txbox(slide, item, 0.85, y, 11.8, 0.55, size=18, color=WHITE)
        y += 0.65

# ── Slide 1: Title ──────────────────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.06, BLUE)
rect(s, 0, 7.44, 13.33, 0.06, VIOLET)
txbox(s, "SkillSetu AI", 1, 1.6, 11, 1.2, size=60, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "AI Workforce Intelligence Platform", 1, 2.9, 11, 0.7, size=28, color=BLUE, align=PP_ALIGN.CENTER)
txbox(s, "Karnataka Government Initiative  •  Hackathon Prototype", 1, 3.65, 11, 0.5, size=18, color=MUTED, align=PP_ALIGN.CENTER)
txbox(s, "ಕರ್ನಾಟಕದ ಕಾರ್ಮಿಕರಿಗಾಗಿ AI ಮೌಲ್ಯಮಾಪನ ವ್ಯವಸ್ಥೆ", 1, 4.25, 11, 0.6, size=20, color=AMBER, align=PP_ALIGN.CENTER)
# stat boxes
for i, (val, lbl) in enumerate([("3", "Languages"), ("10+", "Skill Categories"), ("30", "Districts")]):
    x = 2.5 + i * 3
    rect(s, x, 5.2, 2.4, 1.2, RGBColor(0x1E, 0x29, 0x4A), line_color=BLUE)
    txbox(s, val, x, 5.25, 2.4, 0.65, size=32, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    txbox(s, lbl, x, 5.85, 2.4, 0.45, size=14, color=MUTED, align=PP_ALIGN.CENTER)

# ── Slide 2: Problem Statement ───────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, RED)
txbox(s, "The Problem", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, RED)

problems = [
    ("Karnataka has 30+ districts with millions of skilled workers — but no unified assessment system.", RED),
    ("Manual interviews are slow, inconsistent, and prone to fraud.", AMBER),
    ("Language barriers prevent fair evaluation of Kannada & Hindi-speaking workers.", AMBER),
    ("Government agencies lack real-time workforce readiness data.", BLUE),
    ("No automated way to detect duplicate or fraudulent candidates.", RED),
]
y = 1.5
for text, color in problems:
    rect(s, 0.55, y + 0.14, 0.1, 0.1, color)
    txbox(s, text, 0.85, y, 11.8, 0.55, size=17, color=WHITE)
    y += 0.7

# ── Slide 3: Solution Overview ───────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, BLUE)
txbox(s, "Our Solution", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, BLUE)
txbox(s, "An end-to-end AI-powered platform that automates workforce assessment for Karnataka's skilled workers.",
      0.4, 1.3, 12.5, 0.7, size=18, color=MUTED)

cards = [
    ("AI Video Interviews", "Automated interviews with real-time transcription & emotion analysis", BLUE),
    ("Smart Scoring", "Multi-dimensional: communication, confidence, relevance, authenticity", VIOLET),
    ("Fraud Detection", "Detects duplicates, multiple faces, audio inconsistencies automatically", RED),
    ("Multilingual", "Kannada-first with full Hindi & English support", AMBER),
    ("Analytics Dashboard", "District-wise insights & workforce readiness reports", EMERALD),
    ("Auto Classification", "Job Ready / Needs Training / Manual Review", CYAN),
]
cols = 3
for i, (title, desc, color) in enumerate(cards):
    col = i % cols
    row = i // cols
    x = 0.4 + col * 4.3
    y = 2.3 + row * 2.2
    rect(s, x, y, 4.0, 1.9, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    rect(s, x + 0.15, y + 0.18, 0.18, 0.18, color)
    txbox(s, title, x + 0.45, y + 0.12, 3.4, 0.45, size=15, bold=True, color=WHITE)
    txbox(s, desc,  x + 0.15, y + 0.65, 3.7, 0.9,  size=12, color=MUTED)

# ── Slide 4: How It Works ────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, VIOLET)
txbox(s, "How It Works", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, VIOLET)

steps = [
    ("01", "Register", "Candidate fills basic details & selects preferred language (Kannada / Hindi / English)", BLUE),
    ("02", "AI Interview", "Answers 5 questions via webcam. AI transcribes speech in real-time.", VIOLET),
    ("03", "AI Analysis", "Scores communication, confidence, relevance, authenticity & detects fraud", AMBER),
    ("04", "Classification", "Auto-classified into workforce readiness categories with full report", EMERALD),
]
for i, (num, title, desc, color) in enumerate(steps):
    x = 0.5 + i * 3.2
    rect(s, x, 1.6, 2.9, 4.8, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    txbox(s, num,   x, 1.7, 2.9, 1.0, size=40, bold=True, color=RGBColor(0x1E, 0x29, 0x4A), align=PP_ALIGN.CENTER)
    txbox(s, title, x, 2.7, 2.9, 0.6, size=18, bold=True, color=color, align=PP_ALIGN.CENTER)
    txbox(s, desc,  x + 0.15, 3.4, 2.6, 2.5, size=13, color=MUTED)
    if i < 3:
        txbox(s, "→", x + 2.85, 3.5, 0.4, 0.5, size=22, bold=True, color=color)

# ── Slide 5: AI Scoring Engine ───────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, AMBER)
txbox(s, "AI Scoring Engine", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, AMBER)

metrics = [
    ("Communication", "Clarity, vocabulary, sentence structure", BLUE, 82),
    ("Confidence",    "Voice tone, pace, hesitation patterns",   VIOLET, 74),
    ("Relevance",     "Answer alignment to the question asked",  EMERALD, 88),
    ("Authenticity",  "Inverse of fraud score — genuine response", AMBER, 91),
    ("Fraud Score",   "Multiple faces, audio mismatch, visibility", RED, 15),
]
y = 1.5
for name, desc, color, sample in metrics:
    txbox(s, name, 0.5, y, 3.0, 0.45, size=16, bold=True, color=WHITE)
    txbox(s, desc, 0.5, y + 0.38, 3.5, 0.35, size=12, color=MUTED)
    # bar background
    rect(s, 4.0, y + 0.1, 7.5, 0.28, RGBColor(0x1E, 0x29, 0x4A))
    # bar fill
    rect(s, 4.0, y + 0.1, 7.5 * sample / 100, 0.28, color)
    txbox(s, f"{sample}%", 11.7, y + 0.05, 0.9, 0.4, size=14, bold=True, color=color)
    y += 1.0

# ── Slide 6: Fraud Detection ─────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, RED)
txbox(s, "Fraud Detection System", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, RED)

flags = [
    ("Multiple Faces",     "Detects when more than one person is visible in the frame", "HIGH",   RED),
    ("Audio Mismatch",     "Voice pattern inconsistency — coached or pre-recorded audio", "MEDIUM", AMBER),
    ("Low Visibility",     "Face partially obscured or camera covered during interview", "LOW",    BLUE),
    ("Duplicate Candidate","Same person re-registering under different details", "HIGH",   RED),
]
y = 1.5
for flag, desc, sev, color in flags:
    rect(s, 0.4, y, 12.2, 1.0, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    rect(s, 0.4, y, 0.25, 1.0, color)
    txbox(s, flag, 0.8, y + 0.08, 4.5, 0.45, size=16, bold=True, color=WHITE)
    txbox(s, desc, 0.8, y + 0.52, 8.5, 0.38, size=13, color=MUTED)
    rect(s, 10.5, y + 0.28, 1.8, 0.38, color)
    txbox(s, sev, 10.5, y + 0.28, 1.8, 0.38, size=13, bold=True, color=DARK_BG, align=PP_ALIGN.CENTER)
    y += 1.2

# ── Slide 7: Admin Dashboard ─────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, EMERALD)
txbox(s, "Admin Dashboard", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, EMERALD)

stats = [
    ("20",  "Total Candidates", BLUE),
    ("6",   "Job Ready",        EMERALD),
    ("2",   "Fraud Flags",      RED),
    ("4",   "Pending Reviews",  AMBER),
]
for i, (val, lbl, color) in enumerate(stats):
    x = 0.4 + i * 3.2
    rect(s, x, 1.5, 2.9, 1.5, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    txbox(s, val, x, 1.55, 2.9, 0.9, size=40, bold=True, color=color, align=PP_ALIGN.CENTER)
    txbox(s, lbl, x, 2.4, 2.9, 0.45, size=14, color=MUTED, align=PP_ALIGN.CENTER)

txbox(s, "Dashboard Modules", 0.4, 3.3, 12.5, 0.5, size=18, bold=True, color=WHITE)
modules = [
    ("Candidate Table", "Search, filter, sort all candidates with status badges", BLUE),
    ("Charts & Analytics", "Weekly trends, district breakdown, language & classification pie charts", VIOLET),
    ("Fraud Review Panel", "Dedicated view for high-risk candidates with flag details", RED),
    ("Reports Export", "Download workforce readiness reports by district or skill", EMERALD),
]
y = 3.9
for i, (mod, desc, color) in enumerate(modules):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    yy = y + row * 1.4
    rect(s, x, yy, 6.1, 1.2, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    txbox(s, mod,  x + 0.2, yy + 0.08, 5.7, 0.45, size=15, bold=True, color=color)
    txbox(s, desc, x + 0.2, yy + 0.55, 5.7, 0.55, size=12, color=MUTED)

# ── Slide 8: Tech Stack ──────────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, CYAN)
txbox(s, "Technology Stack", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, CYAN)

tech = [
    ("Frontend",    ["Next.js 14 (App Router)", "TypeScript", "Tailwind CSS", "Framer Motion", "Recharts"], BLUE),
    ("AI / ML",     ["Speech-to-Text (multilingual)", "Emotion & Confidence Analysis", "Fraud Detection CV", "NLP Relevance Scoring"], VIOLET),
    ("State & Data",["Zustand (global store)", "Mock data layer (scalable to real API)", "REST-ready architecture"], EMERALD),
    ("Deployment",  ["Vercel (frontend)", "AWS Lambda (AI inference)", "AWS S3 (video storage)", "AWS RDS (candidate DB)"], AMBER),
]
for i, (cat, items, color) in enumerate(tech):
    col = i % 2
    row = i // 2
    x = 0.4 + col * 6.5
    y = 1.5 + row * 2.8
    rect(s, x, y, 6.1, 2.5, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    txbox(s, cat, x + 0.2, y + 0.1, 5.7, 0.5, size=17, bold=True, color=color)
    for j, item in enumerate(items):
        rect(s, x + 0.3, y + 0.7 + j * 0.42, 0.1, 0.1, color)
        txbox(s, item, x + 0.55, y + 0.65 + j * 0.42, 5.3, 0.38, size=13, color=WHITE)

# ── Slide 9: Impact & Scale ──────────────────────────────────────────────────
s = add_slide()
bg(s)
accent_bar(s, EMERALD)
txbox(s, "Impact & Scale", 0.4, 0.3, 12.5, 0.8, size=32, bold=True, color=WHITE)
rect(s, 0.4, 1.15, 2.5, 0.04, EMERALD)

impacts = [
    ("30 Districts", "Full coverage across Karnataka", BLUE),
    ("3 Languages", "Kannada, Hindi, English — inclusive assessment", AMBER),
    ("10+ Skill Trades", "Construction to Healthcare Aide", VIOLET),
    ("Real-time Results", "Instant classification — no waiting weeks", EMERALD),
    ("Fraud-proof", "AI flags suspicious candidates automatically", RED),
    ("Scalable", "Cloud-native — handles thousands of interviews/day", CYAN),
]
for i, (title, desc, color) in enumerate(impacts):
    col = i % 3
    row = i // 3
    x = 0.4 + col * 4.3
    y = 1.5 + row * 2.5
    rect(s, x, y, 4.0, 2.2, RGBColor(0x0F, 0x17, 0x2A), line_color=color)
    rect(s, x + 0.15, y + 0.2, 0.25, 0.25, color)
    txbox(s, title, x + 0.55, y + 0.15, 3.3, 0.55, size=17, bold=True, color=color)
    txbox(s, desc,  x + 0.15, y + 0.85, 3.7, 1.1,  size=13, color=MUTED)

# ── Slide 10: Thank You / CTA ────────────────────────────────────────────────
s = add_slide()
bg(s)
rect(s, 0, 0, 13.33, 0.06, BLUE)
rect(s, 0, 7.44, 13.33, 0.06, VIOLET)
txbox(s, "SkillSetu AI", 1, 1.5, 11, 1.0, size=52, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txbox(s, "Bridging Karnataka's Skilled Workforce with Opportunity", 1, 2.6, 11, 0.7, size=22, color=BLUE, align=PP_ALIGN.CENTER)
txbox(s, "ಕರ್ನಾಟಕದ ಕಾರ್ಮಿಕರಿಗಾಗಿ AI ಮೌಲ್ಯಮಾಪನ ವ್ಯವಸ್ಥೆ", 1, 3.4, 11, 0.6, size=20, color=AMBER, align=PP_ALIGN.CENTER)

for i, (label, url, color) in enumerate([
    ("Begin Interview", "/interview", BLUE),
    ("Admin Dashboard", "/admin",     VIOLET),
]):
    x = 3.2 + i * 4.0
    rect(s, x, 4.3, 3.2, 0.85, color)
    txbox(s, label, x, 4.3, 3.2, 0.85, size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txbox(s, "© 2026 Karnataka Government Workforce Initiative  •  Hackathon Prototype",
      1, 6.6, 11, 0.5, size=13, color=MUTED, align=PP_ALIGN.CENTER)

prs.save("SkillSetu_AI_Presentation.pptx")
print("Saved: SkillSetu_AI_Presentation.pptx")
